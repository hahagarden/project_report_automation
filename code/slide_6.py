from pptx import Presentation
from pptx.util import Inches

prs=Presentation()

img_path='/git/project_report_automation/flower.jpg'

blank_slide_layout = prs.slide_layouts[6] #6: 제목/내용이 없는 빈 슬라이드
slide=prs.slides.add_slide(blank_slide_layout)

left=top=Inches(1)
width=height=Inches(1)
#width, height가 없을 경우 원본 사이즈로
pic=slide.shapes.add_picture(img_path, left, top, width=width, height=height)

left=Inches(3)
width=Inches(5.5)
height=Inches(4)
pic=slide.shapes.add_picture(img_path, left, top, width=width, height=height)

prs.save('test2.pptx')

