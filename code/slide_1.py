from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

bullet_slide_layout = prs.slide_layouts[1]  # 1:제목 및 내용 슬라이드
slide = prs.slides.add_slide(bullet_slide_layout)  # 기존에 있던 슬라이드에 추가

# 제목
title_shape = slide.placeholders[0]
title_shape.text = 'Adding a Bullet Slide'

# 내용
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text = 'Find the bullet slide layout'

# 단락추가
p = tf.add_paragraph()
p.text = 'Use _TextFrame.text for first bullet'
p.level = 1  # 1:들여쓰기 레벨

# 단락추가
p = tf.add_paragraph()
p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
p.level = 2  # 2:들여쓰기 레벨

prs.save('test.pptx')
