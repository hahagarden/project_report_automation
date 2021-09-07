from pptx import Presentation #라이브러리
from pptx.util import Inches #사진,표 등을 그리기 위해

prs=Presentation()
for i in range(0,11):
    print(f"-------{i}-------")
    slide=prs.slides.add_slide(prs.slide_layouts[i])
    for shape in slide.placeholders:
        print('%d %s' %(shape.placeholder_format.idx, shape.name))


        